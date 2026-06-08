"""装备合成系统 设计汇报 PPT — premium-soft 主题"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ===== 设计系统: premium-soft =====
BG      = (255, 251, 245)   # 暖奶油
HEADING = (45, 36, 24)      # 深咖啡
BODY    = (45, 36, 24)      # 正文
ACCENT  = (196, 149, 106)   # 铜金
MUTED   = (122, 110, 94)    # 灰褐

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== Helper functions =====
def set_bg(slide, rgb=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*rgb)

def add_title_slide(title, subtitle, date="v1.0", team=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # 装饰线
    line = slide.shapes.add_shape(1, Inches(5.5), Inches(2.8), Inches(2.3), Inches(0))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*ACCENT)
    line.line.fill.background()
    # 标题
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(44); p.font.bold = True
    p.font.color.rgb = RGBColor(*HEADING); p.alignment = PP_ALIGN.CENTER
    # 副标题
    p2 = tf.add_paragraph(); p2.text = subtitle
    p2.font.size = Pt(24); p2.font.color.rgb = RGBColor(*MUTED); p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)
    # 底部信息
    tb2 = slide.shapes.add_textbox(Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.6))
    tf2 = tb2.text_frame
    p3 = tf2.paragraphs[0]; p3.text = f"{date}   |   {team}"; p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(*MUTED); p3.alignment = PP_ALIGN.CENTER
    return slide

def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # 标题
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(32); p.font.bold = True
    p.font.color.rgb = RGBColor(*ACCENT)
    # 标题下划线
    line = slide.shapes.add_shape(1, Inches(1), Inches(1.35), Inches(3.5), Inches(0))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(*ACCENT); line.line.fill.background()
    # 要点
    tb2 = slide.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(10.7), Inches(5))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"•  {b}"; p.font.size = Pt(18); p.font.color.rgb = RGBColor(*BODY)
        p.space_after = Pt(14); p.line_spacing = Pt(27)
    return slide

def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(32); p.font.bold = True
    p.font.color.rgb = RGBColor(*ACCENT)
    # 下划线
    line = slide.shapes.add_shape(1, Inches(1), Inches(1.35), Inches(3.5), Inches(0))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(*ACCENT); line.line.fill.background()
    # 表格
    n_rows = len(rows) + 1; n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(1), Inches(1.9), Inches(11.3), Inches(5)).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j); cell.text = h
        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(14); p.font.bold = True
        p.font.color.rgb = RGBColor(*HEADING)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(245, 240, 232)  # accentLt
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j); cell.text = str(val)
            p = cell.text_frame.paragraphs[0]; p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(*BODY)
            if i % 2 == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(250, 246, 240)
    return slide

# ====== 7 Slides ======

add_title_slide(
    "装备合成系统 · 龙鳞工坊",
    "收集材料 → 合成装备 → 强化成长 → 挑战更强敌人",
    date="v1.0 · 2025-07-16",
    team="游戏设计部"
)

add_bullet_slide("项目概述", [
    "核心循环：战斗/探索 → 收集龙鳞·龙心·矿石 → 合成工坊合成装备 → 强化提升 → 挑战更强敌人 → 获得稀有材料",
    "目标受众：18-35岁，喜欢RPG成长感+收集养成，偏好奇幻题材，手机/PC双端",
    "差异化定位：配方确定产物类型（龙鳞×5+龙心×1→龙鳞剑），属性在580~650间随机——兼顾确定性和惊喜感",
    "竞品对比：原神圣遗物随机过强 vs 怪物猎人锻造过于确定 → 本作取中，保留'手造'仪式感",
    "稀有材料可锁定1条属性范围，提高下限避免极品变废品",
])

add_bullet_slide("系统设计", [
    "战斗系统：碰撞盒检测+韧性条硬直机制（破韧2s），0.2s精确弹反窗口返还50%伤害，火>草>水>火元素克制",
    "养成体系：武器Lv.80/防具Lv.60/饰品Lv.50封顶，费用阶梯指数增长 Cost=100×1.15^Level",
    "合成工坊：预览区→投放材料到3个槽位→查看预期属性范围→确认合成，材料不足时按钮灰态",
    "装备属性随机：攻击力580~650 / 暴击率8~15% / 暴伤20~30%，稀有材料可锁定1条属性下限",
    "辅助系统：合成图鉴自动记录已解锁配方，筛选排序快速查找装备，容量50/50时底栏变红",
])

add_table_slide("交互与UI设计 — 界面清单 (9/9 对应原型图)", 
    ["序号", "界面名称", "优先级", "关键说明"],
    [
        ["01", "装备背包主界面", "P0", "3列网格+Tab栏+底栏合成/分解入口"],
        ["02", "装备详情面板", "P0", "3D模型旋转+主属性进度条+副属性列表"],
        ["03", "合成入口/配方列表", "P0", "搜索框+分类Tab+配方卡片(充足/不足/未解锁)"],
        ["04", "合成工坊", "P0", "预览区+3材料槽位+预期属性范围+开始合成"],
        ["05", "合成结果弹窗", "P0", "暗色遮罩+光柱升起+属性展示+查看详情/装备"],
        ["06", "装备强化界面", "P1", "属性对比+消耗材料+成功率100%+强化按钮"],
        ["07", "材料分解界面", "P1", "多选装备+预计获得汇总+分解按钮"],
        ["08", "合成图鉴", "P2", "已解锁8/24+已解锁/未解锁Tab+解锁条件"],
        ["09", "筛选/排序弹窗", "P2", "类型/星级checkbox+排序radio+重置/确定"],
    ]
)

add_bullet_slide("数值设计", [
    "伤害公式：最终伤害=(攻击力×技能倍率−防御力×0.6)×(1+暴击伤害加成)×元素克制系数",
    "生存模型：EHP=HP÷(1−减伤率)，普通怪TTK 3~8s，BOSS TTK 60~120s",
    "经济系统：金币(战斗掉落+分解)→合成强化 / 钻石(充值+成就)→扩充背包 / 合成精华(分解高品质)→锁定属性",
    "强化消耗：武器Cost=100×1.15^Level / 防具Cost=80×1.12^Level / 饰品Cost=60×1.10^Level",
    "敌人成长：1-5波1.0倍→6-10波1.5倍(闪避格挡)→16-20波4.0倍(AOE狂暴)→21-25波6.5倍(双形态元素免疫)",
])

add_bullet_slide("美术需求", [
    "视觉风格：高奢柔软 Premium Soft — 暖奶油底#FFFBF5+深咖啡#2D2418+铜金#C4956A，双重边框嵌套卡片",
    "主角动画7组：Idle(循环60帧)/Run(24帧三段)/Jump(1s)/轻攻击1-2-3(各0.4s)/重攻击(1s蓄力+释放)/Hit(0.33s)/Death(1.33s)",
    "特效清单P0：合成光柱金色800ms+星级点亮stagger150ms+材料放入绿色脉冲",
    "特效清单P1-P2：强化成功金色上扬飞字+分解金粉飞出+按钮hover上浮+GodRays斜射光柱",
    "场景分层：前景UI/粒子(1.0×)·中景可玩区域(1.0×)·远景1建筑山体(0.5×)·远景2天空云(0.2×)",
])

add_bullet_slide("总结与下一步", [
    "核心优势：材料确定产物类型提供目标感，属性范围随机提供每次合成的惊喜——玩家享受'打造'而非纯抽奖",
    "风险缓解：属性范围过宽→稀有材料锁定下限；背包容量50→监控使用率；暗色调下对比度≥4.5:1",
    "扩展方向：PVP排行榜展示最优装备、赛季限定配方、公会材料交易所、装备外观幻化系统",
    "下一步：④ game-asset-requirements 梳理美术/音频/特效资源清单 → ⑤ 整合为HTML汇报文档",
])

# ===== Save =====
desktop = os.path.expanduser("~/Desktop")
output = os.path.join(desktop, "装备合成系统_设计汇报_PPT_test.pptx")
prs.save(output)
print(f"✅ PPT 已保存: {output}")
print(f"📊 7 页 | 📐 16:9 | 🎨 premium-soft")
